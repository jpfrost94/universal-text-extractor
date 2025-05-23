import os
import mimetypes
import tempfile
import re
import csv
import json
from defusedxml import ElementTree as ET
import email
from io import StringIO
from datetime import datetime
import importlib

# Importing potentially needed libraries
try:
    import pandas as pd
    has_pandas = True
except ImportError:
    has_pandas = False

try:
    import fitz  # PyMuPDF
    has_pymupdf = True
except ImportError:
    has_pymupdf = False

try:
    import PyPDF2
    has_pypdf2 = True
except ImportError:
    has_pypdf2 = False

try:
    from PIL import Image
    has_pil = True
except ImportError:
    has_pil = False

try:
    import docx2txt
    has_docx2txt = True
except ImportError:
    has_docx2txt = False

try:
    from docx import Document
    has_python_docx = True
except ImportError:
    has_python_docx = False

try:
    from pptx import Presentation
    has_python_pptx = True
except ImportError:
    has_python_pptx = False

try:
    from bs4 import BeautifulSoup
    has_bs4 = True
except ImportError:
    has_bs4 = False

try:
    from striprtf.striprtf import rtf_to_text
    has_striprtf = True
except ImportError:
    has_striprtf = None

try:
    from odf import text as odf_text, teletype as odf_teletype
    from odf.opendocument import load as odf_load
    from odf.table import Table, TableRow, TableCell
    has_odf = True
except ImportError:
    has_odf = False

try:
    import ebooklib
    from ebooklib import epub
    has_ebooklib = True
except ImportError:
    has_ebooklib = False

try:
    import extract_msg
    has_extract_msg = True
except ImportError:
    has_extract_msg = False

# Import OCR utils
from .ocr_utils import perform_ocr
from .image_processing import preprocess_image

# Define supported file types by category
SUPPORTED_FILE_TYPES = {
    "Documents": ["pdf", "doc", "docx", "rtf", "odt", "txt"],
    "Images": ["jpg", "jpeg", "png", "tiff", "tif", "bmp", "gif", "webp", "heic", "heif"],
    "Presentations": ["ppt", "pptx", "odp"],
    "Spreadsheets": ["xls", "xlsx", "ods", "csv"],
    "Web": ["html", "htm", "xml"],
    "Email": ["eml", "msg"],
    "Ebooks": ["epub"]
}

# List of supported image formats for easier lookup
SUPPORTED_IMAGE_FORMATS = ["jpg", "jpeg", "png", "tiff", "tif", "bmp", "gif", "webp", "heic", "heif"]

def detect_file_type(file_path):
    """
    Detect the type/format of a file based on extension and content.
    
    Args:
        file_path: Path to the file
    
    Returns:
        String indicating the detected file type
    """
    # Get file extension
    _, ext = os.path.splitext(file_path)
    ext = ext.lower().lstrip('.')
    
    # Check extension against supported types
    for category, extensions in SUPPORTED_FILE_TYPES.items():
        if ext in extensions:
            return ext
    
    # Try to determine content type if extension not recognized
    content_type, _ = mimetypes.guess_type(file_path)
    
    if content_type:
        # Map mime type to a general category
        if content_type.startswith('image/'):
            return "image"
        elif content_type.startswith('text/'):
            return "text"
        elif 'pdf' in content_type:
            return "pdf"
        elif 'msword' in content_type or 'vnd.openxmlformats-officedocument.wordprocessingml' in content_type:
            return "doc"
        elif 'vnd.ms-powerpoint' in content_type or 'vnd.openxmlformats-officedocument.presentationml' in content_type:
            return "ppt"
        elif 'vnd.ms-excel' in content_type or 'vnd.openxmlformats-officedocument.spreadsheetml' in content_type:
            return "xls"
    
    # Default case
    return "unknown"

def extract_text_from_file(file_path, file_type=None, use_ocr=False, ocr_language="eng", preprocess_params=None):
    """
    Extract text from a file based on its type.
    
    Args:
        file_path: Path to the file
        file_type: Type of the file (if known, otherwise will be auto-detected)
        use_ocr: Whether to use OCR for images and PDFs
        ocr_language: Language for OCR
        preprocess_params: Parameters for image preprocessing
    
    Returns:
        Tuple of (extracted text, whether OCR was used, log messages)
    """
    log_messages = []
    ocr_used = False
    
    # Auto-detect file type if not provided
    if not file_type:
        file_type = detect_file_type(file_path)
        log_messages.append(f"Auto-detected file type: {file_type}")
    
    # Convert file type to lowercase for case-insensitive comparison
    file_type = file_type.lower()
    
    extracted_text = ""
    
    try:
        # Extract text based on file type
        if file_type == "pdf":
            extracted_text, ocr_used, pdf_logs = extract_text_from_pdf(file_path, use_ocr, ocr_language, preprocess_params)
            log_messages.extend(pdf_logs)
        
        elif file_type in SUPPORTED_IMAGE_FORMATS:
            if use_ocr:
                # Process image for OCR if requested
                if preprocess_params and preprocess_params.get("enhance", False):
                    log_messages.append("Preprocessing image before OCR")
                    image = preprocess_image(file_path, preprocess_params)
                    extracted_text = perform_ocr(image, ocr_language)
                else:
                    extracted_text = perform_ocr(file_path, ocr_language)
                ocr_used = True
                log_messages.append(f"Performed OCR on image with language: {ocr_language}")
            else:
                log_messages.append("OCR not enabled for image file")
                if has_pil:
                    # Basic image metadata without OCR
                    img = Image.open(file_path)
                    width, height = img.size
                    extracted_text = f"[Image: {width}x{height}, Format: {img.format}]\n"
                    extracted_text += "OCR was not enabled. Enable OCR to extract text content from this image."
                else:
                    extracted_text = "[Image file - OCR not enabled]"
        
        elif file_type in ["doc", "docx"]:
            extracted_text, doc_logs = extract_text_from_doc(file_path)
            log_messages.extend(doc_logs)
        
        elif file_type in ["ppt", "pptx"]:
            extracted_text, ppt_logs = extract_text_from_ppt(file_path)
            log_messages.extend(ppt_logs)
        
        elif file_type in ["xls", "xlsx", "csv", "ods"]:
            extracted_text, sheet_logs = extract_text_from_spreadsheet(file_path, file_type)
            log_messages.extend(sheet_logs)
        
        elif file_type == "txt":
            with open(file_path, 'r', encoding='utf-8', errors='replace') as file:
                extracted_text = file.read()
            log_messages.append("Extracted text from TXT file")
        
        elif file_type == "rtf":
            extracted_text, rtf_logs = extract_text_from_rtf(file_path)
            log_messages.extend(rtf_logs)
        
        elif file_type in ["html", "htm"]:
            extracted_text, html_logs = extract_text_from_html(file_path)
            log_messages.extend(html_logs)
        
        elif file_type == "xml":
            extracted_text, xml_logs = extract_text_from_xml(file_path)
            log_messages.extend(xml_logs)
        
        elif file_type in ["odt", "odp"]:
            extracted_text, odf_logs = extract_text_from_odf(file_path)
            log_messages.extend(odf_logs)
        
        elif file_type == "epub":
            extracted_text, epub_logs = extract_text_from_epub(file_path)
            log_messages.extend(epub_logs)
        
        elif file_type in ["eml", "msg"]:
            extracted_text, email_logs = extract_text_from_email(file_path, file_type)
            log_messages.extend(email_logs)
        
        else:
            # Fallback for unsupported or unknown file types
            log_messages.append(f"Unsupported file type: {file_type}")
            extracted_text = f"[Unsupported file type: {file_type}]"
    
    except Exception as e:
        log_messages.append(f"Error extracting text: {str(e)}")
        extracted_text = f"[Error extracting text: {str(e)}]"
    
    return extracted_text, ocr_used, log_messages

def extract_text_from_pdf(pdf_path, use_ocr=False, ocr_language="eng", preprocess_params=None):
    """
    Extract text from a PDF file.
    
    Args:
        pdf_path: Path to the PDF file
        use_ocr: Whether to use OCR
        ocr_language: Language for OCR
        preprocess_params: Parameters for image preprocessing
    
    Returns:
        Tuple of (extracted text, whether OCR was used, log messages)
    """
    logs = []
    extracted_text = ""
    ocr_used = False
    
    # Try PyMuPDF first
    if has_pymupdf:
        try:
            logs.append("Attempting to extract text using PyMuPDF")
            extracted_text = ""
            doc = fitz.open(pdf_path)
            
            for page_num in range(len(doc)):
                page = doc[page_num]
                page_text = page.get_text()
                
                # If page has no text and OCR is enabled, try OCR
                if not page_text.strip() and use_ocr:
                    logs.append(f"Page {page_num+1} appears to be scanned/image-based. Attempting OCR.")
                    pix = page.get_pixmap()
                    
                    # Save pixmap to a temporary image file
                    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp_img:
                        pix.save(tmp_img.name)
                        tmp_img_path = tmp_img.name
                    
                    # Perform OCR on the image
                    if preprocess_params and preprocess_params.get("enhance", False):
                        # Preprocess the image before OCR
                        img = preprocess_image(tmp_img_path, preprocess_params)
                        page_text = perform_ocr(img, ocr_language)
                    else:
                        page_text = perform_ocr(tmp_img_path, ocr_language)
                    
                    # Remove the temporary image file
                    os.unlink(tmp_img_path)
                    ocr_used = True
                
                extracted_text += f"\n--- Page {page_num+1} ---\n{page_text}\n"
            
            doc.close()
            
            if not extracted_text.strip():
                logs.append("No text extracted with PyMuPDF")
            else:
                logs.append(f"Successfully extracted text from PDF using PyMuPDF")
                return extracted_text, ocr_used, logs
                
        except Exception as e:
            logs.append(f"PyMuPDF extraction failed: {str(e)}")
    else:
        logs.append("PyMuPDF not available, trying alternative methods")
    
    # Fallback to PyPDF2
    if has_pypdf2 and not extracted_text:
        try:
            logs.append("Attempting to extract text using PyPDF2")
            extracted_text = ""
            with open(pdf_path, 'rb') as file:
                reader = PyPDF2.PdfReader(file)
                
                for page_num in range(len(reader.pages)):
                    page = reader.pages[page_num]
                    page_text = page.extract_text()
                    
                    # If page has no text and OCR is enabled, we would need to convert to image
                    # But PyPDF2 doesn't have built-in image conversion
                    
                    extracted_text += f"\n--- Page {page_num+1} ---\n{page_text}\n"
            
            if not extracted_text.strip():
                logs.append("No text extracted with PyPDF2")
            else:
                logs.append(f"Successfully extracted text from PDF using PyPDF2")
                return extracted_text, ocr_used, logs
                
        except Exception as e:
            logs.append(f"PyPDF2 extraction failed: {str(e)}")
    
    # If no text has been extracted yet and OCR is enabled, try full PDF OCR
    if not extracted_text.strip() and use_ocr:
        # Try to convert PDF to images and OCR
        logs.append("PDF appears to be scanned/image-based. Attempting full document OCR.")
        
        try:
            from pdf2image import convert_from_path
            
            extracted_text = ""
            # Convert PDF to images
            images = convert_from_path(pdf_path)
            
            for i, img in enumerate(images):
                # Save image to a temporary file
                with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp_img:
                    img.save(tmp_img.name)
                    tmp_img_path = tmp_img.name
                
                # Perform OCR on the image
                if preprocess_params and preprocess_params.get("enhance", False):
                    # Preprocess the image before OCR
                    processed_img = preprocess_image(tmp_img_path, preprocess_params)
                    page_text = perform_ocr(processed_img, ocr_language)
                else:
                    page_text = perform_ocr(tmp_img_path, ocr_language)
                
                # Remove the temporary image file
                os.unlink(tmp_img_path)
                
                extracted_text += f"\n--- Page {i+1} ---\n{page_text}\n"
                ocr_used = True
            
            logs.append(f"Completed full document OCR using pdf2image and OCR")
            
        except ImportError:
            logs.append("pdf2image not available for PDF to image conversion")
        except Exception as e:
            logs.append(f"Full document OCR failed: {str(e)}")
    
    # Final fallback message if nothing worked
    if not extracted_text.strip():
        extracted_text = "[No text could be extracted from this PDF. It may be scanned, image-based, or protected.]"
        logs.append("Could not extract any text from the PDF")
    
    return extracted_text, ocr_used, logs

def extract_text_from_doc(doc_path):
    """
    Extract text from a DOC/DOCX file.
    
    Args:
        doc_path: Path to the DOC/DOCX file
    
    Returns:
        Tuple of (extracted text, log messages)
    """
    logs = []
    extracted_text = ""
    
    # Try docx2txt first (for DOCX files)
    if has_docx2txt and doc_path.lower().endswith(".docx"):
        try:
            logs.append("Attempting to extract text using docx2txt")
            extracted_text = docx2txt.process(doc_path)
            
            if extracted_text.strip():
                logs.append("Successfully extracted text using docx2txt")
                return extracted_text, logs
            else:
                logs.append("No text extracted with docx2txt")
        except Exception as e:
            logs.append(f"docx2txt extraction failed: {str(e)}")
    
    # Try python-docx as an alternative
    if has_python_docx and doc_path.lower().endswith(".docx"):
        try:
            logs.append("Attempting to extract text using python-docx")
            doc = Document(doc_path)
            
            # Extract paragraphs
            paragraphs_text = [paragraph.text for paragraph in doc.paragraphs]
            
            # Extract tables
            tables_text = []
            for table in doc.tables:
                for row in table.rows:
                    row_text = [cell.text for cell in row.cells]
                    tables_text.append(" | ".join(row_text))
            
            # Combine all text
            extracted_text = "\n".join(paragraphs_text)
            
            if tables_text:
                extracted_text += "\n\n--- Tables ---\n"
                extracted_text += "\n".join(tables_text)
            
            if extracted_text.strip():
                logs.append("Successfully extracted text using python-docx")
                return extracted_text, logs
            else:
                logs.append("No text extracted with python-docx")
        except Exception as e:
            logs.append(f"python-docx extraction failed: {str(e)}")
    
    # For DOC files or if other methods fail, indicate limitation
    if doc_path.lower().endswith(".doc"):
        extracted_text = "[DOC (legacy format) file detected. For best results, convert to DOCX format.]"
        logs.append("Legacy DOC format detected, limited extraction support")
    
    # If we got here and have no text, it means all methods failed
    if not extracted_text:
        extracted_text = "[Could not extract text from this document. It may be protected or corrupted.]"
        logs.append("All document extraction methods failed")
    
    return extracted_text, logs

def extract_text_from_ppt(ppt_path):
    """
    Extract text from a PPT/PPTX file.
    
    Args:
        ppt_path: Path to the PPT/PPTX file
    
    Returns:
        Tuple of (extracted text, log messages)
    """
    logs = []
    extracted_text = ""
    
    # Check if python-pptx is available and the file is PPTX
    if has_python_pptx and ppt_path.lower().endswith(".pptx"):
        try:
            logs.append("Attempting to extract text using python-pptx")
            prs = Presentation(ppt_path)
            
            slide_texts = []
            for slide_idx, slide in enumerate(prs.slides, 1):
                slide_text = [f"--- Slide {slide_idx} ---"]
                
                # Get slide title if available
                if slide.shapes.title:
                    slide_text.append(f"Title: {slide.shapes.title.text}")
                
                # Extract text from all shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        # Skip if it's the title we already added
                        if shape == slide.shapes.title:
                            continue
                        slide_text.append(shape.text)
                
                slide_texts.append("\n".join(slide_text))
            
            extracted_text = "\n\n".join(slide_texts)
            
            if extracted_text.strip():
                logs.append("Successfully extracted text using python-pptx")
            else:
                logs.append("No text found in presentation")
                extracted_text = "[No text found in this presentation. It may contain only images or other non-text elements.]"
                
        except Exception as e:
            logs.append(f"python-pptx extraction failed: {str(e)}")
            extracted_text = f"[Error extracting presentation text: {str(e)}]"
    else:
        if ppt_path.lower().endswith(".ppt"):
            extracted_text = "[PPT (legacy format) file detected. For best results, convert to PPTX format.]"
            logs.append("Legacy PPT format detected, limited extraction support")
        else:
            extracted_text = "[Could not extract text from this presentation. Required library not available.]"
            logs.append("python-pptx library not available")
    
    return extracted_text, logs

def extract_text_from_spreadsheet(file_path, file_type):
    """
    Extract text from a spreadsheet file (XLS, XLSX, CSV, ODS).
    
    Args:
        file_path: Path to the spreadsheet file
        file_type: Type of spreadsheet file
    
    Returns:
        Tuple of (extracted text, log messages)
    """
    logs = []
    extracted_text = ""
    
    # Handle CSV files directly
    if file_type == "csv":
        try:
            logs.append("Processing CSV file")
            with open(file_path, 'r', encoding='utf-8', errors='replace') as file:
                reader = csv.reader(file)
                rows = list(reader)
                
                # Format as a table-like structure
                for row in rows:
                    extracted_text += " | ".join(row) + "\n"
            
            logs.append("Successfully extracted data from CSV")
            return extracted_text, logs
        except Exception as e:
            logs.append(f"CSV extraction failed: {str(e)}")
    
    # For Excel files (XLSX, XLS) or ODS, try using pandas
    if has_pandas and file_type in ["xlsx", "xls", "ods"]:
        try:
            logs.append(f"Attempting to extract data from {file_type.upper()} using pandas")
            
            # Read all sheets
            xl = pd.ExcelFile(file_path)
            sheet_names = xl.sheet_names
            
            sheet_texts = []
            for sheet_name in sheet_names:
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                
                # Convert dataframe to string representation
                output = StringIO()
                df.to_csv(output, sep="|", index=False)
                sheet_text = f"--- Sheet: {sheet_name} ---\n{output.getvalue()}"
                sheet_texts.append(sheet_text)
            
            extracted_text = "\n\n".join(sheet_texts)
            logs.append(f"Successfully extracted data from {len(sheet_names)} sheets")
            return extracted_text, logs
        
        except Exception as e:
            logs.append(f"Pandas extraction failed: {str(e)}")
    
    # For ODS, try using ODF if pandas failed
    if has_odf and file_type == "ods" and not extracted_text:
        try:
            logs.append("Attempting to extract data from ODS using odfpy")
            
            doc = odf_load(file_path)
            tables = doc.getElementsByType(Table)
            
            table_texts = []
            for table_idx, table in enumerate(tables, 1):
                table_text = [f"--- Table {table_idx} ---"]
                
                # Process table rows
                for row in table.getElementsByType(TableRow):
                    cells = row.getElementsByType(TableCell)
                    row_values = []
                    
                    for cell in cells:
                        # Get all text nodes
                        text_nodes = cell.getElementsByType(odf_text.P)
                        cell_text = " ".join([odf_teletype.extractText(node) for node in text_nodes])
                        row_values.append(cell_text or "")
                    
                    table_text.append(" | ".join(row_values))
                
                table_texts.append("\n".join(table_text))
            
            extracted_text = "\n\n".join(table_texts)
            logs.append(f"Successfully extracted data from {len(tables)} tables using odfpy")
            return extracted_text, logs
        
        except Exception as e:
            logs.append(f"ODF extraction failed: {str(e)}")
    
    # If we reach here, all methods failed
    if not extracted_text:
        if file_type in ["xlsx", "xls", "ods"]:
            format_upper = file_type.upper()
            extracted_text = f"[Could not extract data from {format_upper} file. Required libraries not available or file is corrupted/protected.]"
        else:
            extracted_text = "[Unsupported spreadsheet format or extraction failed.]"
        
        logs.append("All spreadsheet extraction methods failed")
    
    return extracted_text, logs

def extract_text_from_rtf(rtf_path):
    """
    Extract text from an RTF file.
    
    Args:
        rtf_path: Path to the RTF file
    
    Returns:
        Tuple of (extracted text, log messages)
    """
    logs = []
    
    if has_striprtf:
        try:
            logs.append("Attempting to extract text from RTF using striprtf")
            with open(rtf_path, 'r', encoding='utf-8', errors='replace') as file:
                rtf_content = file.read()
            
            extracted_text = rtf_to_text(rtf_content)
            logs.append("Successfully extracted text from RTF")
            return extracted_text, logs
        
        except Exception as e:
            logs.append(f"RTF extraction failed: {str(e)}")
    else:
        logs.append("striprtf library not available")
    
    # Fallback to basic file reading if striprtf failed or isn't available
    try:
        logs.append("Attempting fallback RTF extraction")
        with open(rtf_path, 'r', encoding='utf-8', errors='replace') as file:
            rtf_content = file.read()
        
        # Very basic RTF tag removal (not comprehensive)
        # Remove RTF control words and groups
        text_without_controls = re.sub(r'\\[a-z0-9]+(-?[0-9]+)?[ ]?|\{|\}|\\', ' ', rtf_content)
        # Remove multiple spaces and other basic cleanup
        cleaned_text = re.sub(r'\s+', ' ', text_without_controls).strip()
        
        logs.append("Used fallback method for RTF extraction")
        return cleaned_text, logs
    
    except Exception as e:
        logs.append(f"Fallback RTF extraction failed: {str(e)}")
        return f"[Error extracting RTF text: {str(e)}]", logs

def extract_text_from_html(html_path):
    """
    Extract text from an HTML file.
    
    Args:
        html_path: Path to the HTML file
    
    Returns:
        Tuple of (extracted text, log messages)
    """
    logs = []
    
    try:
        with open(html_path, 'r', encoding='utf-8', errors='replace') as file:
            html_content = file.read()
        
        if has_bs4:
            logs.append("Extracting text from HTML using BeautifulSoup")
            soup = BeautifulSoup(html_content, 'html.parser')
            
            # Extract title
            title = soup.title.string if soup.title else "Untitled"
            
            # Extract text
            # Remove script and style elements
            for script in soup(["script", "style"]):
                script.extract()
            
            # Get text
            text = soup.get_text(separator='\n')
            
            # Clean up whitespace
            lines = (line.strip() for line in text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            text = '\n'.join(chunk for chunk in chunks if chunk)
            
            extracted_text = f"Title: {title}\n\n{text}"
            logs.append("Successfully extracted text from HTML")
            
        else:
            logs.append("BeautifulSoup not available, using simple HTML tag removal")
            # Simple HTML tag removal (not as good as BeautifulSoup)
            text = re.sub(r'<[^>]+>', ' ', html_content)
            text = re.sub(r'\s+', ' ', text).strip()
            extracted_text = text
            logs.append("Used basic regex for HTML text extraction")
    
    except Exception as e:
        logs.append(f"HTML extraction failed: {str(e)}")
        extracted_text = f"[Error extracting HTML text: {str(e)}]"
    
    return extracted_text, logs

def extract_text_from_xml(xml_path):
    """
    Extract text from an XML file.
    
    Args:
        xml_path: Path to the XML file
    
    Returns:
        Tuple of (extracted text, log messages)
    """
    logs = []
    
    try:
        logs.append("Extracting text from XML")
        
        # Parse XML file
        tree = ET.parse(xml_path)
        root = tree.getroot()
        
        # Function to recursively extract text from elements
        def extract_text_from_element(element, indent=0):
            result = []
            indent_str = "  " * indent
            
            # Element tag and attributes
            attribs = "" if not element.attrib else " " + " ".join([f"{k}='{v}'" for k, v in element.attrib.items()])
            result.append(f"{indent_str}<{element.tag}{attribs}>")
            
            # Element text content
            if element.text and element.text.strip():
                result.append(f"{indent_str}  {element.text.strip()}")
            
            # Process child elements
            for child in element:
                result.extend(extract_text_from_element(child, indent + 1))
            
            # Element closing tag
            result.append(f"{indent_str}</{element.tag}>")
            
            return result
        
        # Extract all text
        text_parts = extract_text_from_element(root)
        extracted_text = "\n".join(text_parts)
        
        logs.append("Successfully extracted text from XML")
    
    except Exception as e:
        logs.append(f"XML extraction failed: {str(e)}")
        
        # Fallback to basic tag extraction
        try:
            with open(xml_path, 'r', encoding='utf-8', errors='replace') as file:
                xml_content = file.read()
            
            # Simple regex to extract text between tags
            text_parts = re.findall(r'>([^<]+)<', xml_content)
            extracted_text = "\n".join([part.strip() for part in text_parts if part.strip()])
            
            logs.append("Used fallback method for XML extraction")
            
        except Exception as fallback_e:
            logs.append(f"Fallback XML extraction also failed: {str(fallback_e)}")
            extracted_text = f"[Error extracting XML text: {str(e)}]"
    
    return extracted_text, logs

def extract_text_from_odf(odf_path):
    """
    Extract text from an ODF file (ODT, ODP).
    
    Args:
        odf_path: Path to the ODF file
    
    Returns:
        Tuple of (extracted text, log messages)
    """
    logs = []
    
    if has_odf:
        try:
            logs.append("Extracting text from ODF document")
            doc = odf_load(odf_path)
            
            # Extract paragraphs
            paragraphs = doc.getElementsByType(odf_text.P)
            text_parts = [odf_teletype.extractText(paragraph) for paragraph in paragraphs]
            
            # Extract headings (if any)
            headings = []
            for i in range(1, 6):  # H1 to H5
                h_type = getattr(odf_text, f"H{i}", None)
                if h_type:
                    h_elements = doc.getElementsByType(h_type)
                    headings.extend([f"Heading {i}: {odf_teletype.extractText(h)}" for h in h_elements])
            
            # Combine all text
            all_parts = headings + text_parts
            extracted_text = "\n".join([part for part in all_parts if part.strip()])
            
            logs.append("Successfully extracted text from ODF document")
            
        except Exception as e:
            logs.append(f"ODF extraction failed: {str(e)}")
            extracted_text = f"[Error extracting ODF text: {str(e)}]"
    else:
        logs.append("odfpy library not available")
        extracted_text = "[Could not extract text from ODF document. Required library not available.]"
    
    return extracted_text, logs

def extract_text_from_epub(epub_path):
    """
    Extract text from an EPUB file.
    
    Args:
        epub_path: Path to the EPUB file
    
    Returns:
        Tuple of (extracted text, log messages)
    """
    logs = []
    
    if has_ebooklib:
        try:
            logs.append("Extracting text from EPUB")
            book = epub.read_epub(epub_path)
            
            # Extract metadata
            title = book.get_metadata('DC', 'title')[0][0] if book.get_metadata('DC', 'title') else "Untitled"
            creator = book.get_metadata('DC', 'creator')[0][0] if book.get_metadata('DC', 'creator') else "Unknown"
            
            metadata = f"Title: {title}\nAuthor: {creator}\n\n"
            
            # Extract content
            contents = []
            for item in book.get_items():
                if item.get_type() == ebooklib.ITEM_DOCUMENT:
                    # Get HTML content
                    html_content = item.get_content().decode('utf-8')
                    
                    # Parse with BeautifulSoup if available
                    if has_bs4:
                        soup = BeautifulSoup(html_content, 'html.parser')
                        # Remove script and style elements
                        for script in soup(["script", "style"]):
                            script.extract()
                        text = soup.get_text(separator='\n')
                    else:
                        # Simple HTML tag removal as fallback
                        text = re.sub(r'<[^>]+>', ' ', html_content)
                        text = re.sub(r'\s+', ' ', text).strip()
                    
                    contents.append(text)
            
            extracted_text = metadata + "\n".join(contents)
            logs.append("Successfully extracted text from EPUB")
            
        except Exception as e:
            logs.append(f"EPUB extraction failed: {str(e)}")
            extracted_text = f"[Error extracting EPUB text: {str(e)}]"
    else:
        logs.append("ebooklib library not available")
        extracted_text = "[Could not extract text from EPUB. Required library not available.]"
    
    return extracted_text, logs

def extract_text_from_email(email_path, email_type):
    """
    Extract text from an email file (EML, MSG).
    
    Args:
        email_path: Path to the email file
        email_type: Type of email file ('eml' or 'msg')
    
    Returns:
        Tuple of (extracted text, log messages)
    """
    logs = []
    
    if email_type == 'eml':
        try:
            logs.append("Extracting text from EML file")
            with open(email_path, 'r', encoding='utf-8', errors='replace') as file:
                eml_content = file.read()
            
            # Parse email
            msg = email.message_from_string(eml_content)
            
            # Extract headers
            from_address = msg.get('From', 'Unknown')
            to_address = msg.get('To', 'Unknown')
            subject = msg.get('Subject', 'No Subject')
            date = msg.get('Date', 'Unknown')
            
            headers = f"From: {from_address}\nTo: {to_address}\nSubject: {subject}\nDate: {date}\n\n"
            
            # Extract body
            body = ""
            if msg.is_multipart():
                for part in msg.walk():
                    content_type = part.get_content_type()
                    content_disp = str(part.get('Content-Disposition'))
                    
                    # Skip attachments
                    if 'attachment' in content_disp:
                        continue
                    
                    # Get text content
                    if content_type == 'text/plain':
                        body_part = part.get_payload(decode=True)
                        charset = part.get_content_charset()
                        if charset:
                            try:
                                body += body_part.decode(charset)
                            except UnicodeDecodeError:
                                body += body_part.decode('utf-8', errors='replace')
                        else:
                            body += body_part.decode('utf-8', errors='replace')
                    
                    elif content_type == 'text/html':
                        # If no plain text is found, use HTML
                        if not body:
                            html_part = part.get_payload(decode=True)
                            charset = part.get_content_charset()
                            
                            if charset:
                                try:
                                    html_content = html_part.decode(charset)
                                except UnicodeDecodeError:
                                    html_content = html_part.decode('utf-8', errors='replace')
                            else:
                                html_content = html_part.decode('utf-8', errors='replace')
                            
                            # Parse HTML if BeautifulSoup is available
                            if has_bs4:
                                soup = BeautifulSoup(html_content, 'html.parser')
                                body += soup.get_text(separator='\n')
                            else:
                                # Basic HTML tag removal
                                body += re.sub(r'<[^>]+>', ' ', html_content)
            else:
                # Not multipart - get the payload directly
                body = msg.get_payload(decode=True).decode('utf-8', errors='replace')
            
            extracted_text = headers + body
            logs.append("Successfully extracted text from EML file")
            
        except Exception as e:
            logs.append(f"EML extraction failed: {str(e)}")
            extracted_text = f"[Error extracting EML text: {str(e)}]"
    
    elif email_type == 'msg' and has_extract_msg:
        try:
            logs.append("Extracting text from MSG file")
            msg = extract_msg.Message(email_path)
            
            # Extract headers
            from_address = msg.sender
            to_address = msg.to
            subject = msg.subject
            date = msg.date
            
            headers = f"From: {from_address}\nTo: {to_address}\nSubject: {subject}\nDate: {date}\n\n"
            
            # Get body
            body = msg.body
            
            extracted_text = headers + body
            logs.append("Successfully extracted text from MSG file")
            
        except Exception as e:
            logs.append(f"MSG extraction failed: {str(e)}")
            extracted_text = f"[Error extracting MSG text: {str(e)}]"
    
    else:
        if email_type == 'msg' and not has_extract_msg:
            logs.append("extract_msg library not available")
            extracted_text = "[Could not extract text from MSG file. Required library not available.]"
        else:
            logs.append(f"Unsupported email type: {email_type}")
            extracted_text = f"[Unsupported email type: {email_type}]"
    
    return extracted_text, logs
